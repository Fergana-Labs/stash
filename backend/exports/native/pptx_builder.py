"""list[SlideSpec] → PPTX bytes via python-pptx native shapes.

Public entry: `build_pptx(specs, html, image_fetcher=None)`. `html` is
passed through to `hybrid_raster.rasterize_targets` for any raster
fallbacks the probe asked for. `image_fetcher` is optional so tests can
inject a stub instead of touching the network.
"""

from __future__ import annotations

import io
import logging

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from ..constants import (
    SLIDE_HEIGHT_EMU,
    SLIDE_HEIGHT_PX,
    SLIDE_WIDTH_EMU,
    SLIDE_WIDTH_PX,
)
from .hybrid_raster import rasterize_targets
from .image_fetch import ImageFetcher
from .spec import BBox, Paragraph, ShapeSpec, SlideSpec

logger = logging.getLogger(__name__)


def px_to_emu_x(px: float) -> int:
    """Convert CSS px to PPTX EMU along the x-axis. Clamped to the slide
    rect so a stray out-of-bounds element doesn't produce negative EMU."""
    return max(0, min(int(SLIDE_WIDTH_EMU), int(px * int(SLIDE_WIDTH_EMU) / SLIDE_WIDTH_PX)))


def px_to_emu_y(px: float) -> int:
    return max(0, min(int(SLIDE_HEIGHT_EMU), int(px * int(SLIDE_HEIGHT_EMU) / SLIDE_HEIGHT_PX)))


def _bbox_emu(bbox: BBox) -> tuple[int, int, int, int]:
    left = px_to_emu_x(bbox.x)
    top = px_to_emu_y(bbox.y)
    width = max(1, px_to_emu_x(bbox.x2) - left)
    height = max(1, px_to_emu_y(bbox.y2) - top)
    return left, top, width, height


def _hex_to_rgb(color: str | None) -> RGBColor | None:
    if not color or not color.startswith("#") or len(color) != 7:
        return None
    try:
        return RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))
    except ValueError:
        return None


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


async def build_pptx(
    specs: list[SlideSpec],
    html: str,
    image_fetcher: ImageFetcher | None = None,
) -> bytes:
    """Assemble specs into a PPTX. Coordinates Playwright rasterisation
    of background fallbacks and image fetching."""
    fetcher = image_fetcher or ImageFetcher()

    # Pre-resolve every raster target in one Playwright pass.
    raster_targets: list[tuple[int, str]] = []
    for spec in specs:
        if spec.bg_raster_selector:
            raster_targets.append((spec.index, spec.bg_raster_selector))
        for sh in _walk_shapes(spec.shapes):
            if sh.kind == "raster" and sh.raster_selector:
                raster_targets.append((spec.index, sh.raster_selector))
    rasters = await rasterize_targets(html, raster_targets)

    # Pre-fetch every image in parallel-ish (the fetcher caches per-URL,
    # but in v1 we just await sequentially — decks are small).
    for spec in specs:
        for sh in _walk_shapes(spec.shapes):
            if sh.kind == "image" and sh.src:
                await fetcher.fetch(sh.src)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU
    blank_layout = prs.slide_layouts[6]

    for spec in specs:
        slide = prs.slides.add_slide(blank_layout)
        _emit_slide_background(slide, spec, rasters)
        # Sort shapes by z so card backgrounds end up behind their text.
        for sh in sorted(spec.shapes, key=lambda s: s.z):
            await _emit_shape(slide, sh, rasters, fetcher, spec.index)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _walk_shapes(shapes: list[ShapeSpec]):
    for s in shapes:
        yield s
        for row in s.cells:
            for cell in row:
                yield from _walk_shapes([cell])


def _emit_slide_background(slide, spec: SlideSpec, rasters: dict[tuple[int, str], bytes]) -> None:
    if spec.bg_raster_selector:
        png = rasters.get((spec.index, spec.bg_raster_selector))
        if png:
            slide.shapes.add_picture(io.BytesIO(png), 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU)
            return
    rgb = _hex_to_rgb(spec.bg_color)
    if rgb is not None:
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = rgb


async def _emit_shape(slide, sh: ShapeSpec, rasters, fetcher: ImageFetcher, slide_idx: int) -> None:
    if sh.kind == "raster":
        png = rasters.get((slide_idx, sh.raster_selector or ""))
        if png:
            left, top, w, h = _bbox_emu(sh.bbox)
            slide.shapes.add_picture(io.BytesIO(png), left, top, w, h)
        return

    if sh.kind == "image":
        data = await fetcher.fetch(sh.src) if sh.src else None
        if not data:
            return
        left, top, w, h = _bbox_emu(sh.bbox)
        try:
            slide.shapes.add_picture(io.BytesIO(data), left, top, w, h)
        except Exception:
            logger.exception("add_picture failed for %s", sh.src[:120] if sh.src else "?")
        return

    if sh.kind == "table":
        _emit_table(slide, sh)
        return

    # text shape (possibly with a card background)
    _emit_text(slide, sh)


def _emit_text(slide, sh: ShapeSpec) -> None:
    left, top, w, h = _bbox_emu(sh.bbox)
    if sh.bg_color and not sh.paragraphs:
        # Pure card background — rectangle without text.
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        rgb = _hex_to_rgb(sh.bg_color)
        if rgb is not None:
            rect.fill.solid()
            rect.fill.fore_color.rgb = rgb
        rect.line.fill.background()
        return

    if sh.bg_color:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        rgb = _hex_to_rgb(sh.bg_color)
        if rgb is not None:
            rect.fill.solid()
            rect.fill.fore_color.rgb = rgb
        rect.line.fill.background()
        tf = rect.text_frame
    else:
        box = slide.shapes.add_textbox(left, top, w, h)
        tf = box.text_frame

    tf.margin_left = Emu(int(px_to_emu_x(sh.padding_px[3])))
    tf.margin_right = Emu(int(px_to_emu_x(sh.padding_px[1])))
    tf.margin_top = Emu(int(px_to_emu_y(sh.padding_px[0])))
    tf.margin_bottom = Emu(int(px_to_emu_y(sh.padding_px[2])))
    tf.word_wrap = True
    # The bbox came from getBoundingClientRect, which for tight-line-height
    # elements (`<h1 style="line-height:0.95">`) is exactly text-height. PPTX
    # textboxes with default auto_size will visually clip the descender row;
    # tell python-pptx to keep the explicit dimensions but allow visible
    # overflow.
    tf.auto_size = MSO_AUTO_SIZE.NONE

    for i, para in enumerate(sh.paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _populate_paragraph(p, para)


def _populate_paragraph(p, para: Paragraph) -> None:
    p.alignment = _ALIGN_MAP.get(para.align, PP_ALIGN.LEFT)
    # line_height in the spec is a multiplier (1.2 = 120% line height). PPTX
    # paragraph.line_spacing accepts a multiplier directly when it's a float.
    if para.line_height and 0.5 <= para.line_height <= 5.0:
        p.line_spacing = float(para.line_height)
    # python-pptx's default paragraph already has one (empty) run; clear it
    # by replacing the text directly on the first run we add.
    first_run = None
    for run_spec in para.runs:
        text = run_spec.text
        if text == "\n":
            p.add_line_break()
            continue
        r = p.add_run() if first_run is not None else _first_run(p)
        if first_run is None:
            first_run = r
        r.text = text
        # PPTX 13.333" slides hold a 1920 CSS-px viewport in our probe,
        # so effective DPI = 1920 / 13.333 = 144. Converting font px → pt
        # at 96 DPI (the html default) overshoots by 50% and the result
        # is text that's too big for its own bounding box. The slide
        # canvas uses px / 2 = pt instead.
        r.font.size = Pt(max(6, int(run_spec.font_size_px * 72 / 144)))
        r.font.bold = run_spec.bold
        r.font.italic = run_spec.italic
        r.font.underline = run_spec.underline
        if run_spec.font_family:
            r.font.name = run_spec.font_family
        rgb = _hex_to_rgb(run_spec.color)
        if rgb is not None:
            r.font.color.rgb = rgb


def _first_run(p):
    if p.runs:
        return p.runs[0]
    return p.add_run()


def _emit_table(slide, sh: ShapeSpec) -> None:
    rows = sh.cells
    if not rows or not rows[0]:
        return
    n_rows, n_cols = len(rows), max(len(r) for r in rows)
    left, top, w, h = _bbox_emu(sh.bbox)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, w, h)
    table = shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, cell_spec in enumerate(row):
            if c_idx >= n_cols:
                continue
            cell = table.cell(r_idx, c_idx)
            cell.text = ""  # clear default
            if cell_spec.paragraphs:
                tf = cell.text_frame
                for i, para in enumerate(cell_spec.paragraphs):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    _populate_paragraph(p, para)
            if cell_spec.bg_color:
                rgb = _hex_to_rgb(cell_spec.bg_color)
                if rgb is not None:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb
